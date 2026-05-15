#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
生成作概述 PowerPoint 演示文稿
基于图片内容的 PPT 自动生成脚本
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    """创建 PowerPoint 演示文稿"""
    
    # 创建演示文稿对象
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 第一页：作概述
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    # 设置背景色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 248, 255)  # 浅蓝色背景
    
    # 添加标题 "01 作概述"
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(4), Inches(1))
    title_frame = title_box.text_frame
    title_frame.word_wrap = False
    p = title_frame.paragraphs[0]
    p.text = "01  作概述"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)
    
    # 添加蓝色下划线
    line_shape = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(2.8), Inches(0))
    line_shape.line.color.rgb = RGBColor(0, 102, 204)
    line_shape.line.width = Pt(5)
    
    # 添加主要内容文本框
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(5.5), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    # 第一段
    p = text_frame.paragraphs[0]
    p.text = "2026年 是公司战略转型的关键一年，也是我个人职业生涯中具有重要意义的一年。"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.space_after = Pt(10)
    
    # 第二段
    p = text_frame.add_paragraph()
    p.text = "本年度，我担任项目部高级经理职务，主要负责公司重点项目的统筹管理与推进工作。"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.space_after = Pt(10)
    
    # 第三段
    p = text_frame.add_paragraph()
    p.text = "在公司的正确领导和同事们的大力支持下，我紧紧围绕年度工作目标，认真履行岗位职责，顺利完成各项工作任务。"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0, 0, 0)
    
    # 添加四个统计卡片
    cards = [
        ("项目数量", "12个", "全年共负责和参与项目"),
        ("重大项目", "3个", "其中重大项目"),
        ("项目总预算", "3500万元", "项目总预算达"),
        ("项目平均完成率", "98.5%", "较去年同期提升5.2个百分点")
    ]
    
    card_x_start = 0.5
    card_y = 5.8
    card_width = 2.05
    card_height = 1.3
    
    for i, (title, number, desc) in enumerate(cards):
        card_x = card_x_start + (i * 2.2)
        
        # 卡片背景框
        card_shape = slide.shapes.add_shape(1, Inches(card_x), Inches(card_y), Inches(card_width), Inches(card_height))
        card_shape.fill.solid()
        card_shape.fill.fore_color.rgb = RGBColor(245, 250, 255)
        card_shape.line.color.rgb = RGBColor(200, 220, 240)
        card_shape.line.width = Pt(1)
        
        # 卡片标题
        title_box = slide.shapes.add_textbox(Inches(card_x + 0.1), Inches(card_y + 0.1), Inches(card_width - 0.2), Inches(0.35))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        tp = title_tf.paragraphs[0]
        tp.text = title
        tp.font.size = Pt(11)
        tp.font.bold = True
        tp.font.color.rgb = RGBColor(80, 80, 80)
        
        # 数字
        number_box = slide.shapes.add_textbox(Inches(card_x + 0.1), Inches(card_y + 0.45), Inches(card_width - 0.2), Inches(0.5))
        number_tf = number_box.text_frame
        np = number_tf.paragraphs[0]
        np.text = number
        np.font.size = Pt(26)
        np.font.bold = True
        np.font.color.rgb = RGBColor(0, 102, 204)
        
        # 描述文字
        desc_box = slide.shapes.add_textbox(Inches(card_x + 0.1), Inches(card_y + 0.95), Inches(card_width - 0.2), Inches(0.3))
        desc_tf = desc_box.text_frame
        desc_tf.word_wrap = True
        dp = desc_tf.paragraphs[0]
        dp.text = desc
        dp.font.size = Pt(8)
        dp.font.color.rgb = RGBColor(120, 120, 120)
    
    # 添加团队建设内容框
    team_box = slide.shapes.add_textbox(Inches(6.8), Inches(5.3), Inches(3), Inches(1.8))
    team_tf = team_box.text_frame
    team_tf.word_wrap = True
    
    # 团队建设标题
    tp = team_tf.paragraphs[0]
    tp.text = "👥  团队建设成效显著"
    tp.font.size = Pt(15)
    tp.font.bold = True
    tp.font.color.rgb = RGBColor(0, 102, 204)
    tp.space_after = Pt(6)
    
    # 团队建设描述
    tp = team_tf.add_paragraph()
    tp.text = "在团队建设方面，我领导的15人项目团队凝聚力显著增强，团队成员业务能力得到全面提升。"
    tp.font.size = Pt(11)
    tp.font.color.rgb = RGBColor(0, 0, 0)
    tp.space_after = Pt(0)
    
    # 保存文件
    output_file = '作概述.pptx'
    prs.save(output_file)
    print(f"✅ PowerPoint 文件创建成功: {output_file}")
    return output_file

if __name__ == '__main__':
    create_presentation()
